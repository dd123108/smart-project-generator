from flask import Flask, render_template, request, send_file
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import random, io, requests
from duckduckgo_search import DDGS
from PIL import Image
from io import BytesIO
import os

app = Flask(__name__)

TOPICS = [
    "Artificial Intelligence", "Machine Learning", "Cybersecurity",
    "Cloud Computing", "Web Development", "Data Science",
    "Blockchain Technology", "Augmented Reality", "Computer Networks",
    "IoT (Internet of Things)", "Big Data Analytics"
]

PREWRITTEN_SUMMARIES = {
    "Cybersecurity": [
        [
            "Cybersecurity protects systems, networks, and data from cyberattacks.",
            "It includes tools, processes, and practices to guard digital infrastructure.",
            "Key threats include malware, phishing, ransomware, and insider attacks.",
            "Cybersecurity ensures confidentiality, integrity, and availability of information.",
            "With growing digitalization, its importance has become critical for all sectors."
        ],
        [
            "Common cyber threats include viruses, worms, spyware, and trojans.",
            "Phishing uses deceptive emails to steal credentials or deploy malware.",
            "Ransomware locks data until a ransom is paid, disrupting operations.",
            "DDoS attacks overload systems, causing downtime and loss of revenue.",
            "Insider threats involve employees misusing access to harm the organization."
        ],
        [
            "Network security includes firewalls, intrusion detection, and encryption.",
            "Authentication methods like 2FA and biometrics improve access control.",
            "Antivirus software and endpoint protection guard individual devices.",
            "SIEM tools help monitor, detect, and respond to threats in real time.",
            "Regular patching and updates close security loopholes in software systems."
        ],
        [
            "Cybersecurity is vital for banking, healthcare, government, and defense.",
            "In banking, it prevents fraud and protects customer financial data.",
            "Healthcare systems use it to secure patient records and medical devices.",
            "Government agencies guard against espionage and infrastructure attacks.",
            "The defense sector uses advanced systems to counter sophisticated threats."
        ],
        [
            "Key cybersecurity practices include risk assessment and incident response.",
            "Organizations implement policies like BYOD and acceptable use policies.",
            "Training staff reduces human error, a major cause of breaches.",
            "Data backup and recovery plans help in quick restoration post-attacks.",
            "Compliance with standards like ISO 27001 and GDPR is essential."
        ],
        [
            "Ethical hackers help identify and fix security flaws before misuse.",
            "Penetration testing simulates attacks to assess system defenses.",
            "Bug bounty programs reward external researchers for finding vulnerabilities.",
            "Cyber forensics investigates breaches and gathers legal evidence.",
            "AI and machine learning are increasingly used to automate threat detection."
        ],
        [
            "As threats evolve, so must cybersecurity strategies and technologies.",
            "Cloud security and zero-trust architecture are gaining adoption.",
            "The cybersecurity job market is growing rapidly with high demand.",
            "Staying updated with trends and training is crucial for professionals.",
            "Future challenges include quantum threats and AI-powered cyberwarfare."
        ]
    ],
    "Machine Learning": [
        [
            "Machine Learning (ML) enables systems to learn from data without explicit programming.",
            "It is a subset of AI that focuses on building algorithms to improve over time.",
            "ML models identify patterns, make predictions, and automate decisions.",
            "Applications span healthcare, finance, e-commerce, and autonomous systems.",
            "Data quality and quantity are crucial for effective model training."
        ],
        [
            "There are three main types: supervised, unsupervised, and reinforcement learning.",
            "Supervised learning uses labeled data to train models (e.g., regression, classification).",
            "Unsupervised learning discovers hidden patterns in unlabeled data (e.g., clustering).",
            "Reinforcement learning learns optimal actions via rewards and penalties.",
            "Each type serves different tasks and business objectives."
        ],
        [
            "Popular ML algorithms include decision trees, support vector machines, and k-NN.",
            "Neural networks mimic the human brain and power deep learning applications.",
            "Random forests and ensemble methods improve model accuracy.",
            "Naive Bayes is efficient for text classification and spam detection.",
            "Gradient boosting techniques like XGBoost are widely used in competitions."
        ],
        [
            "Model evaluation metrics include accuracy, precision, recall, and F1 score.",
            "Cross-validation ensures models perform well on unseen data.",
            "Overfitting and underfitting are common issues in ML modeling.",
            "Hyperparameter tuning and regularization improve model generalization.",
            "Visualization tools like confusion matrix aid in performance analysis."
        ],
        [
            "Tools and frameworks like Scikit-learn, TensorFlow, and PyTorch simplify ML.",
            "Python is the most popular language for ML development.",
            "Jupyter notebooks help in interactive coding and data visualization.",
            "Cloud platforms like AWS and Google Cloud offer scalable ML services.",
            "MLOps practices help automate ML workflows and deployment."
        ],
        [
            "Ethical concerns include data privacy, bias, and explainability.",
            "Biased training data can lead to unfair or discriminatory models.",
            "Explainable AI aims to make ML decisions transparent to users.",
            "Regulations like GDPR influence how ML systems handle personal data.",
            "Responsible AI development requires accountability and fairness."
        ],
        [
            "ML continues to evolve with trends like AutoML and edge ML.",
            "AutoML automates model selection, training, and deployment.",
            "Edge ML runs models on local devices like phones and IoT sensors.",
            "ML is driving advancements in robotics, NLP, and recommendation systems.",
            "Future innovations will likely combine ML with quantum computing."
        ]
    ],
    "Cloud Computing": [
        ["Cloud computing delivers computing services over the internet on demand.",
         "It includes storage, servers, databases, networking, and software."],
        ["Key models are IaaS, PaaS, and SaaS, offering flexibility to users.",
         "Users can scale resources up or down based on need and cost."],
        ["Major providers include AWS, Microsoft Azure, and Google Cloud Platform.",
         "They offer global infrastructure and security compliance."],
        ["Cloud benefits include cost savings, performance, speed, and efficiency.",
         "Data is accessible from anywhere with internet connectivity."],
        ["Challenges include data security, compliance, and service outages.",
         "Organizations must manage cloud governance and vendor lock-in."],
        ["Hybrid and multi-cloud strategies are becoming popular choices.",
         "They allow balancing between on-premise and cloud services."],
        ["The future includes AI-driven cloud services and serverless computing.",
         "Cloud will be central to digital transformation efforts worldwide."]
    ],
    "Web Development": [
        ["Web development involves building and maintaining websites and web apps.",
         "It includes frontend, backend, and full-stack development."],
        ["Frontend uses HTML, CSS, and JavaScript for user interface.",
         "Frameworks like React, Angular, and Vue enhance frontend performance."],
        ["Backend uses languages like Python, Node.js, and PHP to handle logic.",
         "Databases such as MySQL and MongoDB store and retrieve data."],
        ["Responsive design ensures websites work across devices and screens.",
         "Tools like Bootstrap simplify building responsive layouts."],
        ["Version control with Git helps track and manage code changes.",
         "Platforms like GitHub support collaborative development."],
        ["Web security includes HTTPS, input validation, and authentication.",
         "Developers implement security best practices to protect users."],
        ["Web development trends include JAMstack, PWAs, and serverless apps.",
         "These approaches improve performance and scalability."]
    ],
    "Data Science": [
        ["Data science extracts insights and knowledge from structured and unstructured data.",
         "It blends statistics, programming, and domain expertise."],
        ["The data science process involves data collection, cleaning, and analysis.",
         "Visualization helps in understanding patterns and trends."],
        ["Python and R are popular programming languages in data science.",
         "Libraries like Pandas, NumPy, and Matplotlib are widely used."],
        ["Machine learning plays a key role in predictive analytics.",
         "Data scientists build models to forecast outcomes and automate decisions."],
        ["Big data tools like Hadoop and Spark manage large data sets.",
         "They provide scalability and fast processing capabilities."],
        ["Data ethics focus on responsible data use, privacy, and fairness.",
         "Biased data can lead to misleading or discriminatory results."],
        ["Data science impacts healthcare, finance, marketing, and government.",
         "It drives innovation, efficiency, and better decision-making."]
    ],
    "Blockchain Technology": [
        ["Blockchain is a distributed ledger technology ensuring secure data recording.",
         "It is decentralized, transparent, and tamper-proof."],
        ["Each block contains a list of transactions linked to previous blocks.",
         "This creates an immutable chain that is hard to alter."],
        ["Bitcoin introduced blockchain for digital currency.",
         "Now, it is used beyond crypto in supply chain, voting, and identity."],
        ["Smart contracts automate agreements without intermediaries.",
         "Ethereum is a leading platform supporting smart contracts."],
        ["Consensus mechanisms like Proof of Work and Proof of Stake validate transactions.",
         "They secure the network and maintain integrity."],
        ["Challenges include scalability, energy use, and regulation.",
         "New solutions aim to address these limitations."],
        ["Blockchain’s future lies in Web3, DeFi, and cross-industry adoption.",
         "It promises trustless and efficient digital ecosystems."]
    ],
    "Augmented Reality": [
        ["Augmented Reality (AR) overlays digital content on the real world.",
         "It enhances user experience by blending virtual elements with reality."],
        ["AR is used in gaming, retail, education, and industrial applications.",
         "Popular examples include Pokémon GO and IKEA Place."],
        ["AR devices include smartphones, smart glasses, and headsets.",
         "They use sensors, cameras, and processors to map surroundings."],
        ["AR development tools include ARKit, ARCore, and Unity.",
         "These platforms enable developers to create immersive experiences."],
        ["Challenges include privacy, user safety, and hardware limitations.",
         "AR systems must avoid distractions and ensure data protection."],
        ["AR in education offers interactive and engaging learning environments.",
         "In retail, it enables virtual try-ons and product visualization."],
        ["The future of AR includes integration with AI and 5G.",
         "This will lead to more intelligent, seamless, and real-time interactions."]
    ],
    "Computer Networks": [
        ["Computer networks connect devices to share resources and data.",
         "Types include LAN, WAN, MAN, and PAN."],
        ["Networking hardware includes routers, switches, and access points.",
         "They manage data routing, distribution, and access."],
        ["Protocols like TCP/IP, HTTP, and FTP govern communication.",
         "They ensure reliable data transfer and connectivity."],
        ["Network topologies define structure—star, bus, ring, and mesh.",
         "Each has pros and cons based on use case."],
        ["Security involves firewalls, encryption, and intrusion detection.",
         "VPNs protect data over public networks."],
        ["Monitoring tools track performance and identify issues.",
         "QoS ensures priority handling for critical traffic."],
        ["Future trends include 5G, SDN, and IoT networking.",
         "They promise faster, more flexible, and scalable networks."]
    ],
    "IoT (Internet of Things)": [
        ["IoT connects physical devices to the internet for data exchange.",
         "Examples include smart homes, wearables, and industrial sensors."],
        ["IoT architecture includes devices, gateways, cloud, and applications.",
         "Sensors collect data, which is processed and acted upon remotely."],
        ["Communication protocols include MQTT, CoAP, and Zigbee.",
         "These ensure efficient and secure device interaction."],
        ["IoT platforms help manage devices, data, and integration.",
         "Popular ones are AWS IoT, Azure IoT Hub, and Google IoT Core."],
        ["Security concerns include data breaches, hacking, and privacy risks.",
         "Strong encryption and authentication are essential."],
        ["IoT in healthcare enables remote monitoring and diagnostics.",
         "In industry, it improves automation and predictive maintenance."],
        ["Edge computing brings processing closer to devices for faster response.",
         "AI and 5G will further enhance IoT capabilities."]
    ],
    "Big Data Analytics": [
        ["Big Data refers to massive, diverse, and fast-growing datasets.",
         "It follows the 3Vs: Volume, Velocity, and Variety."],
        ["Analytics turns raw data into actionable insights.",
         "It includes descriptive, predictive, and prescriptive analysis."],
        ["Technologies include Hadoop, Spark, and NoSQL databases.",
         "They handle data storage, processing, and querying."],
        ["Visualization tools like Tableau and Power BI present data effectively.",
         "Clear visuals help in data-driven decision-making."],
        ["Big data use cases span fraud detection, customer analytics, and logistics.",
         "It enables personalization and operational efficiency."],
        ["Data governance ensures accuracy, security, and compliance.",
         "Organizations must manage data quality and access."],
        ["The future includes real-time analytics and integration with AI.",
         "Big data will fuel innovation across all sectors."]
    ],
    "Artificial Intelligence": [
        ["Artificial Intelligence (AI) enables machines to mimic human intelligence.",
         "It includes tasks like learning, reasoning, problem-solving, and language understanding."],
        ["AI applications span across healthcare, finance, and autonomous vehicles.",
         "Examples include voice assistants, recommendation systems, and medical diagnostics."],
        ["AI technologies include machine learning, neural networks, and deep learning.",
         "These enable systems to improve from data and experience."],
        ["Ethical concerns involve bias, job displacement, and data privacy.",
         "Responsible AI development aims to address these issues."],
        ["Narrow AI focuses on specific tasks, while General AI aims for broader capabilities.",
         "General AI is still in theoretical and research phases."],
        ["AI impacts productivity, innovation, and competitive advantage in industries.",
         "Companies invest heavily in AI research and integration."],
        ["Future AI trends include explainable AI, AI governance, and human-AI collaboration.",
         "These aim to make AI more transparent and trustworthy."]
    ]
}

def fetch_summary(topic, num_points):
    topic_clean = topic.strip().lower()
    for predefined_topic, paragraphs in PREWRITTEN_SUMMARIES.items():
        if predefined_topic.lower() == topic_clean:
            return paragraphs[:num_points]
    # fallback to DuckDuckGo (not very detailed)
    with DDGS() as ddgs:
        results = ddgs.text(topic, region='wt-wt', safesearch='Moderate', max_results=1)
        for result in results:
            snippet = result['body']
            words = snippet.split()
            paragraph_length = len(words) // (num_points * 2)
            paragraphs = [
                " ".join(words[i * paragraph_length:(i + 1) * paragraph_length])
                for i in range(num_points * 2)
            ]
            return [paragraphs[i*2:i*2+2] for i in range(num_points)]
    return [["Content not available", "Please try a different topic"]] * num_points

def fetch_images(topic, num_images):
    image_urls = []
    with DDGS() as ddgs:
        results = ddgs.images(topic, max_results=num_images)
        for result in results:
            image_urls.append(result['image'])
    return image_urls

def add_background_image(slide, image_path, prs):
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    slide.shapes.add_picture(image_path, 0, 0, width=slide_width, height=slide_height)

@app.route('/')
def home():
    return render_template('index.html')



@app.route('/generate', methods=['POST'])
def generate():
    name = request.form['name']
    roll = request.form['roll']
    teacher = request.form['teacher']
    topic = request.form['topic']
    num_slides = int(request.form['num_slides'])
    
    if topic == 'random':
        topic = random.choice(TOPICS)

    prs = Presentation()
    slide_layout = prs.slide_layouts[6]

    slide = prs.slides.add_slide(slide_layout)
    if os.path.exists("background.jpg"):
        add_background_image(slide, "background.jpg", prs)

    title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1.5))
    tf = title_box.text_frame
    run = tf.paragraphs[0].add_run()
    run.text = topic
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 51, 102)

    info_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
    tf2 = info_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = f"Submitted by: {name}\nRoll No: {roll}\nSubmitted to: {teacher}"
    p2.font.size = Pt(32)
    p2.font.color.rgb = RGBColor(80, 80, 80)

    summaries = fetch_summary(topic, num_slides)
    images = fetch_images(topic, num_slides)

    for i in range(num_slides):
        slide = prs.slides.add_slide(slide_layout)
        if os.path.exists("background.jpg"):
            add_background_image(slide, "background.jpg", prs)

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
        tf = title_box.text_frame
        run = tf.paragraphs[0].add_run()
        run.text = f"{topic} - Slide {i+1}"
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0, 0, 128)

        # Add image above text
        if i < len(images):
            try:
                response = requests.get(images[i], timeout=5)
                image = Image.open(BytesIO(response.content)).convert("RGB")
                image_stream = BytesIO()
                image.save(image_stream, format='PNG')
                image_stream.seek(0)
                slide.shapes.add_picture(image_stream, Inches(0.5), Inches(1), width=Inches(3))
            except:
                pass

        # Bullet content
        content_box = slide.shapes.add_textbox(Inches(4), Inches(1), Inches(5), Inches(5))
        tf = content_box.text_frame
        tf.word_wrap = True
        for bullet in summaries[i]:
            p = tf.add_paragraph()
            p.text = f"• {bullet}"
            p.level = 0
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(0, 0, 0)

    # Conclusion Slide
    slide = prs.slides.add_slide(slide_layout)
    if os.path.exists("background.jpg"):
        add_background_image(slide, "background.jpg", prs)
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1.5))
    tf = title_box.text_frame
    run = tf.paragraphs[0].add_run()
    run.text = "Conclusion"
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 51, 102)

    conclusion_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(3))
    tf2 = conclusion_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = f"{topic} continues to shape the future and holds vast potential for innovation and impact."
    run2.font.size = Pt(32)
    run2.font.color.rgb = RGBColor(0, 0, 0)
    
    run3 = p2.add_run()
    run3.text = f"THANK YOU"
    run3.font.size = Pt(42)
    run3.font.color.rgb = RGBColor(255,0,0)

    ppt_io = io.BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)

    filename = f"{name}_{roll}_{topic.replace(' ', '_')}.pptx"
    return send_file(
        ppt_io,
        as_attachment=True,
        download_name=filename,
        mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
    )

if __name__ == "__main__":  # This is your if statement
    port = int(os.environ.get("PORT", 5000))  # This line is indented to be part of the if block
    app.run(host="0.0.0.0", port=port)  # This line is also indented to be part of the if block
    app.run(debug=True)


